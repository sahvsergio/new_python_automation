
# Reading PowerPoint presentations
from pptx import Presentation

path_to_presentation = 'Función empresarial.pptx'
# presentation object
prs = Presentation(path_to_presentation)
print('Presentation object for Funcion Empresarial:', prs)
print('*'*100)
# slides object
print('Slides are :')
for slide in prs.slides:
    print('slide object:', slide)
    
print('slide has the following objects')
slide1, slide2=prs.slides[0],prs.slides[1]
print('Slides ids:\n',slide1.slide_id,',',slide2.slide_id)
print('Slide Open XML elements:\n',slide1.element,',',slide2.element)
print('Slide Layouts: \n',slide1.slide_layout.name,slide2.slide_layout.name )
print('Shapes in the slides')
i=1
for slide in prs.slides:
    print(slide, i)
    for shape in slide.shapes:
        print('Shape:',shape.shape_type)
    i+=1

text_runs=[]
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
            
print('The text is:',text_runs)
            
            


# Creating and updating presentations, and adding slides
from pptx import Presentation
prs=Presentation()
slide=prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text='Yo, Python'
slide.placeholders[1].text='Yes, it is really awesome'
prs.save('yoPython.pptx')

"""We can also create a new presentation from an existing presentation. 
In the next code example, we take a PowerPoint template, create a new PPT, 
and add a slide to it with text content. We use the following template for this example:"""

from pptx import Presentation 
prs = Presentation('yoPython.pptx') 
 
first_slide = prs.slides[0] 
first_slide.shapes[0].text_frame.paragraphs[0].text = "Hello!" 
 
slide = prs.slides.add_slide(prs.slide_layouts[1]) 
text_frame = slide.shapes[0].text_frame 
p = text_frame.paragraphs[0] 
p.text = "This is a paragraph" 
 
prs.save('new_ppt.pptx')
# Playing with layouts, placeholders, and textboxes
##Adding A two content slide
from pptx import Presentation
prs=Presentation()
two_content_slide_layout=prs.slide_layouts[3]
slide=prs.slides.add_slide(two_content_slide_layout)
shapes=slide.shapes
title_shape=shapes.title
title_shape.text='Adding A two content slide'


body_shape=shapes.placeholders[1]
tf=body_shape.text_frame    
tf.text='This is line 1'

p=tf.add_paragraph()
p.text='Again a line 2'
p.level=1

p=tf.add_paragraph()
p.text='And this is line 3 ...'
p.level=2

prs.save('two_content.pptx')


##adding a textbox
from pptx import Presentation
from pptx.util import Inches, Pt
prs=Presentation()
blank_slide_layout=prs.slide_layouts[6]
slide=prs.slides.add_slide(blank_slide_layout)

txBox=slide.shapes.add_textbox(Inches(2),Inches(2),Inches(5),Inches(1))
tf=txBox.text_frame
tf.text='Wow, I\'m inside a textbox'

p=tf.add_paragraph()
p.text='Adding a new text'
p.font.bold=True
p.font.italic=True
p.font.size=Pt(30)

prs.save('Textbox.pptx')

# Working with different shapes and adding tables
## working with different shapes

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor

prs=Presentation()
title_only_slide_layout=prs.slide_layouts[5]
slide=prs.slides.add_slide(title_only_slide_layout)
shapes=slide.shapes
shapes.title.text='Adding Shapes'


shape1= shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT,Inches(3.5),Inches(2),Inches(2),Inches(2))
shape1.fill.solid()
shape1.fill.fore_color.rgb=RGBColor(0x1E, 0x90, 0xFF)
shape1.fill.fore_color.brightness =0.4

shape1.text= 'See! There is home !'

shape2= shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_HOME,Inches(3.5),Inches(5),Inches(2),Inches(2))
shape.text='Home'
prs.save('shapes.pptx')


##adding tables
from pptx import Presentation
from pptx.util import Inches
prs=Presentation()
title_only_slide_layout=prs.slide_layouts[5]
slide=prs.slides.add_slide(title_only_slide_layout)
shapes=slide.shapes
shapes.title.text='Student Data'
rows=4
cols=3
left=top=Inches(2.0)
width=Inches(6.0)
height=Inches(1.2)

table=shapes.add_table(rows, cols,left,top, width, height).table


table.columns[0].width=Inches(2.0)
table.columns[1].width=Inches(2.0)
table.columns[1].width=Inches(2.0)

table.cell(0,0).text='Sr. No'
table.cell(0,1).text='Student Name'
table.cell(0,2).text='Student Id'


students={
    1:['john',115],
    2:['Mary', 119],
    3:['Alice',101]
    
}

for i in range(len(students)):
    
    table.cell(i+1,0).text=str(i+1)
    table.cell(i+1,1).text=str(students[i+1][0])
    table.cell(i+1,2).text=str(students[i+1][1])
    
prs.save('table.pptx')
# Visual treat with pictures and charts
from pptx import Presentation
from pptx.util import Inches
img_path = 'python-logo.png'
img_path2 = '134559_10150125234465561_1885506_o.jpg'
prs=Presentation()
blank_slide_layout=prs.slide_layouts[6]
slide= prs.slides.add_slide(blank_slide_layout)

left=top=Inches(2)
pic=slide.shapes.add_picture(img_path,left,top,height=Inches(2),width=Inches(3))

left=Inches(2)
top=Inches(5)
height=Inches(2)
pic=slide.shapes.add_picture(img_path2, left, top, height=height)

prs.save('picture.pptx')

##adding a chart
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION

from pptx.util import Inches

prs=Presentation()
slide=prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text='Data Based on regions'

chart_data=ChartData()
chart_data.categories = ['West', 'East',
                         'North', 'South']

chart_data.add_series('Series 1',(0.25,0.35,0.25,0.25,0.15))

x, y, cx, cy=Inches(2), Inches(2),Inches(6),Inches(4.5)

chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE,x,y,cx,cy, chart_data).chart

chart.has_legend=True
chart.legend.position=XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout=False


chart.plots[0].has_data_labels=True
data_labels=chart.plots[0].data_labels
data_labels.number_format='0%'
data_labels.position=XL_LABEL_POSITION.OUTSIDE_END


prs.save('chart.pptx')
# Automating weekly sales reports
